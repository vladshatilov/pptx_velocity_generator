import os
import sys
import traceback
import statistics
from datetime import datetime

from pptx import Presentation
from pptx.chart.data import XyChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor

import io

import openpyxl

try:
    from PIL import Image as _PILImage
    _HAS_PIL = True
except ImportError:
    _HAS_PIL = False

# Formats that python-pptx cannot open directly — converted to PNG in-memory via PIL
_NEEDS_CONVERSION = {'.webp', '.heic', '.avif'}

# ---------------------------------------------------------------------------
# Paths — resolve relative to the executable or script, not cwd
# ---------------------------------------------------------------------------
if getattr(sys, 'frozen', False):
    APP_DIR     = os.path.dirname(sys.executable)
    # onefile builds extract bundled data into _MEIPASS; onedir puts them next to the exe
    _BUNDLE_DIR = getattr(sys, '_MEIPASS', APP_DIR)
else:
    APP_DIR     = os.path.dirname(os.path.abspath(__file__))
    _BUNDLE_DIR = APP_DIR

BASE_FOLDER    = os.path.join(APP_DIR, 'example_folder')
XLSX_PATH      = os.path.join(BASE_FOLDER, 'velocity.xlsx')
OUTPUT_PATH    = os.path.join(APP_DIR, 'velocity.pptx')
ERROR_PATH     = os.path.join(APP_DIR, 'velocity.txt')
_FALLBACK_IMG  = os.path.join(_BUNDLE_DIR, 'fallback.png')  # user-provided placeholder
_GENERATED_PNG = os.path.join(APP_DIR,    'default_image.png')  # auto-generated gray square

# Processed fallback image bytes (white bg removed), loaded once at startup.
# None = not yet initialised; b'' = unavailable/failed; otherwise raw PNG bytes.
_fallback_bytes: bytes = None

SHEET1 = 'Sheet1'
SHEET2 = 'Sheet2'

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------
IMAGE_EXTENSIONS = ('.jpg', '.jpeg', '.gif', '.png', '.bmp',
                    '.tiff', '.tif', '.webp', '.heic', '.avif')

COLUMN_RENAME = {
    'rank':   'Рейтинг',
    'change': 'Изм vs. LY',
    'ros':    'Продажи на ТТ',
    'nd':     'ND, LM',
    'price':  'Микс цена',
}

CATEGORY_COLORS = [
    RGBColor(0x00, 0x6D, 0x50),
    RGBColor(0x44, 0x70, 0xF3),
    RGBColor(0xFF, 0xC0, 0x00),
    RGBColor(0xEB, 0x23, 0x16),
]

FONT         = 'Calibri'
SLIDE_WIDTH  = Cm(29.7)
SLIDE_HEIGHT = Cm(21.0)
DEFAULT_GW   = Cm(20)
DEFAULT_GH   = Cm(15)


# ---------------------------------------------------------------------------
# Utilities
# ---------------------------------------------------------------------------
def uniquify(path):
    base, ext = os.path.splitext(path)
    counter = 1
    while os.path.exists(path):
        path = f"{base}_{counter}{ext}"
        counter += 1
    return path


def write_error(user_message, technical_detail=None):
    lines = [f"ОШИБКА: {user_message}\n"]
    if technical_detail:
        lines.append(f"\nТехническая информация:\n{technical_detail}\n")
    lines.append(f"\nВремя: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
    content = "".join(lines)
    try:
        with open(ERROR_PATH, 'w', encoding='utf-8') as f:
            f.write(content)
    except Exception:
        pass
    print(content)


def get_sku_images_map():
    """Return {sku_name_without_ext: full_path} for all images in BASE_FOLDER."""
    result = {}
    if not os.path.isdir(BASE_FOLDER):
        return result
    for fname in os.listdir(BASE_FOLDER):
        if fname.lower().endswith(IMAGE_EXTENSIONS):
            result[os.path.splitext(fname)[0]] = os.path.join(BASE_FOLDER, fname)
    return result


def get_blank_layout(prs):
    for layout in prs.slide_layouts:
        if layout.name == 'Blank':
            return layout
    return prs.slide_layouts[-1]


def _init_fallback(threshold: int = 230):
    """
    Load fallback.png once at startup, make near-white pixels transparent via PIL,
    and cache the result as raw bytes.  Each call to _place_image that needs the
    fallback creates a fresh BytesIO from these cached bytes — no repeated I/O or
    re-processing.

    threshold: pixels with R, G, B all >= this value are made fully transparent.
    Tune down (e.g. 200) for cream/light-grey backgrounds.
    """
    global _fallback_bytes
    if _fallback_bytes is not None:
        return  # already done

    if not os.path.exists(_FALLBACK_IMG):
        _fallback_bytes = b''
        return

    if _HAS_PIL:
        try:
            with _PILImage.open(_FALLBACK_IMG) as img:
                img = img.convert('RGBA')
                data = img.getdata()
                img.putdata([
                    (r, g, b, 0) if (r >= threshold and g >= threshold and b >= threshold)
                    else (r, g, b, a)
                    for r, g, b, a in data
                ])
                buf = io.BytesIO()
                img.save(buf, format='PNG')
                _fallback_bytes = buf.getvalue()
            return
        except Exception:
            pass

    # PIL unavailable or processing failed — store original bytes as-is
    try:
        with open(_FALLBACK_IMG, 'rb') as f:
            _fallback_bytes = f.read()
    except Exception:
        _fallback_bytes = b''


def _ensure_gray_png():
    """
    Guarantee that _GENERATED_PNG (a plain gray square) exists on disk.
    Called once at startup so it is always available as the last-resort fallback.
    Uses only stdlib — no external dependencies.
    """
    if os.path.exists(_GENERATED_PNG):
        return
    try:
        import struct, zlib
        w, h, gray = 100, 100, 200

        def make_chunk(tag, data):
            payload = tag + data
            return (struct.pack('>I', len(data)) + tag + data
                    + struct.pack('>I', zlib.crc32(payload) & 0xFFFFFFFF))

        raw = b''.join(b'\x00' + bytes([gray] * w) for _ in range(h))
        png = (
            b'\x89PNG\r\n\x1a\n'
            + make_chunk(b'IHDR', struct.pack('>IIBBBBB', w, h, 8, 0, 0, 0, 0))
            + make_chunk(b'IDAT', zlib.compress(raw))
            + make_chunk(b'IEND', b'')
        )
        with open(_GENERATED_PNG, 'wb') as f:
            f.write(png)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Data reading
# ---------------------------------------------------------------------------
def _parse_nd(val):
    if isinstance(val, str):
        s = val.strip()
        if s.endswith('%'):
            return float(s[:-1]) / 100
        return float(s)
    return float(val) if val is not None else 0.0


def read_sheet1():
    wb = openpyxl.load_workbook(XLSX_PATH, read_only=True, data_only=True)

    if SHEET1 not in wb.sheetnames:
        wb.close()
        raise KeyError(
            f"Лист '{SHEET1}' не найден в файле velocity.xlsx. "
            f"Убедитесь, что в файле есть лист с именем Sheet1."
        )

    ws = wb[SHEET1]
    rows_iter = ws.iter_rows(values_only=True)
    header_row = next(rows_iter, None)

    if header_row is None:
        wb.close()
        raise ValueError(
            "Лист Sheet1 пустой — нет строк с данными. "
            "Добавьте заголовки и данные в таблицу на листе Sheet1."
        )

    headers = [str(h).strip().lower() if h is not None else '' for h in header_row]
    required = ['sku', 'nd', 'ros', 'rank', 'change', 'price']
    missing = [c for c in required if c not in headers]
    if missing:
        wb.close()
        raise KeyError(
            f"В таблице не найдены столбцы: {', '.join(missing)}. "
            f"На листе Sheet1 должны быть столбцы: sku, nd, ros, rank, change, price "
            f"(регистр не важен). Проверьте названия заголовков."
        )

    idx = {h: i for i, h in enumerate(headers)}
    data = []

    for row in rows_iter:
        sku_val = row[idx['sku']]
        if sku_val is None:
            continue
        data.append({
            'sku':    str(sku_val).strip(),
            'nd':     _parse_nd(row[idx['nd']]),
            'ros':    float(row[idx['ros']])    if row[idx['ros']]    is not None else 0.0,
            'rank':   int(row[idx['rank']])     if row[idx['rank']]   is not None else 0,
            'change': row[idx['change']],
            'price':  row[idx['price']],
        })

    wb.close()
    return data


def read_sheet2():
    """Return (graph_width, graph_height, thresholds_list)."""
    try:
        wb = openpyxl.load_workbook(XLSX_PATH, read_only=True, data_only=True)
        if SHEET2 not in wb.sheetnames:
            wb.close()
            return DEFAULT_GW, DEFAULT_GH, []

        ws = wb[SHEET2]
        row2 = next(ws.iter_rows(min_row=2, max_row=2, min_col=1, max_col=5, values_only=True), None)
        wb.close()

        if row2 is None:
            return DEFAULT_GW, DEFAULT_GH, []

        a2, b2, c2, d2, e2 = (row2 + (None, None, None, None, None))[:5]
        gw = Cm(float(a2)) if a2 is not None else DEFAULT_GW
        gh = Cm(float(b2)) if b2 is not None else DEFAULT_GH

        thresholds = []
        for val in [c2, d2, e2]:
            if val is not None:
                thresholds.append(int(val))
            else:
                break

        return gw, gh, thresholds

    except Exception:
        return DEFAULT_GW, DEFAULT_GH, []


# ---------------------------------------------------------------------------
# Drawing helpers
# ---------------------------------------------------------------------------
def add_cross(slide, left, top, width, height):
    """Red X — drawn only when every fallback image has failed."""
    for x1, y1, x2, y2 in [
        (left, top, left + width, top + height),
        (left + width, top, left, top + height),
    ]:
        ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
        ln.line.color.rgb = RGBColor(0xFF, 0x00, 0x00)
        ln.line.width = Pt(1.5)


def draw_border(slide, width, height):
    """Rectangular border starting at 0,0."""
    for x1, y1, x2, y2 in [
        (0, 0, width, 0),
        (width, 0, width, height),
        (width, height, 0, height),
        (0, height, 0, 0),
    ]:
        slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)


def add_median_crosshair(slide, chart_frame, median_nd, median_ros, max_nd, max_ros):
    """
    Draws red horizontal + vertical lines at the median position.
    The plot-area offsets are heuristic fractions of the chart bounding box.
    """
    if max_nd <= 0 or max_ros <= 0:
        return

    left_frac, right_frac = 0.14, 0.03
    top_frac, bottom_frac = 0.05, 0.14

    cl, ct = chart_frame.left, chart_frame.top
    cw, ch = chart_frame.width, chart_frame.height

    pl = cl + int(cw * left_frac)
    pt = ct + int(ch * top_frac)
    pw = int(cw * (1 - left_frac - right_frac))
    ph = int(ch * (1 - top_frac - bottom_frac))

    y_med = pt + ph - int(median_ros / max_ros * ph)
    h = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, pl, y_med, pl + pw, y_med)
    h.line.color.rgb = RGBColor(0xE8, 0x00, 0x00)
    h.line.width = Pt(1.5)

    x_med = pl + int(median_nd / max_nd * pw)
    v = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x_med, pt, x_med, pt + ph)
    v.line.color.rgb = RGBColor(0xE8, 0x00, 0x00)
    v.line.width = Pt(1.5)


def _get_native_size(path):
    """Return (width_px, height_px) for an image file, or None if PIL is unavailable."""
    if not _HAS_PIL:
        return None
    try:
        with _PILImage.open(path) as img:
            return img.size  # (width, height)
    except Exception:
        return None


def _to_pptx_src(path):
    """
    Return a source suitable for slide.shapes.add_picture().
    python-pptx cannot open WebP/HEIC/AVIF directly, so when PIL is available
    those formats are converted to PNG in a BytesIO buffer on the fly.
    For all other formats the original path string is returned unchanged.
    """
    if not path:
        return path
    if os.path.splitext(path)[1].lower() in _NEEDS_CONVERSION and _HAS_PIL:
        try:
            buf = io.BytesIO()
            with _PILImage.open(path) as img:
                img.save(buf, format='PNG')
            buf.seek(0)
            return buf
        except Exception:
            pass
    return path


def _place_image(slide, sku_path, left, top, height, cross_w, rotate_tall=False):
    """
    Place an image preserving aspect ratio (only `height` is constrained).

    Fallback chain (tried in order until one succeeds):
      1. SKU-specific image  (sku_path)
      2. fallback.png        (_FALLBACK_IMG — user-provided, bundled in exe)
      3. gray PNG square     (_GENERATED_PNG — auto-created at startup)
      4. red cross           (last resort)

    rotate_tall — when True and PIL is available, images with native H > 2×W are
    rotated 90° CCW.  The stored position is adjusted so that the visual top-left
    stays at (left, top) after rotation.

    Returns True when the SKU-specific image was placed successfully.
    """
    has_sku = bool(sku_path and os.path.exists(sku_path))

    def _try(src, path_for_size=None):
        """
        Attempt to place an image.  src can be a file path (str) or a BytesIO.
        path_for_size is used only for _get_native_size (BytesIO has no path).
        Returns True on success.
        """
        if src is None:
            return False
        if isinstance(src, str):
            if not os.path.exists(src):
                return False
            path_for_size = src
            src = _to_pptx_src(src)   # convert WebP/HEIC/AVIF → PNG if needed
        try:
            native = _get_native_size(path_for_size) if (rotate_tall and path_for_size) else None
            if native and native[1] > 2 * native[0]:       # H > 2×W → rotate
                w_n, h_n = native
                auto_w = height * w_n // h_n
                delta  = (height - auto_w) // 2
                pic = slide.shapes.add_picture(
                    src, int(left + delta), int(top - delta), height=int(height)
                )
                pic.rotation = 270                          # 90° counter-clockwise
            else:
                slide.shapes.add_picture(src, int(left), int(top), height=int(height))
            return True
        except Exception:
            return False

    # 1. SKU-specific image
    if has_sku and _try(sku_path):
        return True

    # 2. Fallback image (white bg removed, cached in _fallback_bytes)
    #    Fresh BytesIO each call so position is always at 0.
    fb = io.BytesIO(_fallback_bytes) if _fallback_bytes else None
    if _try(fb) or _try(_GENERATED_PNG):
        return False    # fallback used — SKU still reported as missing

    # 3. Nothing worked
    add_cross(slide, int(left), int(top), int(cross_w), int(height))
    return False


def _set_cell_text(cell, text, font_size=Pt(8), bold=False, center=False):
    tf = cell.text_frame
    tf.clear()
    tf.word_wrap = False
    if center:
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    if center:
        p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = str(text) if text is not None else ''
    run.font.name = FONT
    run.font.size = font_size
    run.font.bold = bold


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def build_slide_one(prs, sku_data):
    """Slide 1: XY scatter at 75% slide size, centered, with median annotation."""
    slide = prs.slides.add_slide(get_blank_layout(prs))

    chart_data = XyChartData()
    for rec in sku_data:
        s = chart_data.add_series(str(rec['sku']))
        s.add_data_point(rec['nd'], rec['ros'])

    chart_w    = int(SLIDE_WIDTH  * 0.75)
    chart_h    = int(SLIDE_HEIGHT * 0.75)
    chart_left = (SLIDE_WIDTH  - chart_w) // 2
    chart_top  = (SLIDE_HEIGHT - chart_h) // 2

    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.XY_SCATTER, chart_left, chart_top, chart_w, chart_h, chart_data
    )

    chart = chart_frame.chart
    chart.value_axis.minimum_scale    = 0.0
    chart.category_axis.minimum_scale = 0.0

    nd_vals  = [r['nd']  for r in sku_data]
    ros_vals = [r['ros'] for r in sku_data]
    med_nd   = statistics.median(nd_vals)
    med_ros  = statistics.median(ros_vals)

    add_median_crosshair(
        slide, chart_frame,
        med_nd, med_ros,
        max(nd_vals), max(ros_vals),
    )

    # Footnote: explain the red crosshair
    note_w = int(Cm(16))
    note_h = int(Cm(0.6))
    note_x = int(SLIDE_WIDTH  - note_w - Cm(0.4))
    note_y = int(SLIDE_HEIGHT - note_h - Cm(0.3))
    txb = slide.shapes.add_textbox(note_x, note_y, note_w, note_h)
    p = txb.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = "Пересечение красных линий — медиана"
    run.font.name = FONT
    run.font.size = Pt(8)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0xE8, 0x00, 0x00)


def build_slide_two(prs, sku_data, graph_w, graph_h, sku_images_map):
    """Slide 2: images placed by nd/ros coordinates inside a border frame.
    Aspect ratio preserved (height ≤ 2 cm). No rotation on this slide."""
    slide = prs.slides.add_slide(get_blank_layout(prs))
    draw_border(slide, graph_w, graph_h)

    nd_vals  = [r['nd']  for r in sku_data]
    ros_vals = [r['ros'] for r in sku_data]
    max_nd   = max(nd_vals) if nd_vals else 1
    max_ros  = max(ros_vals) if ros_vals else 1
    img_h    = Cm(2)
    not_found = set()

    for rec in sku_data:
        sku      = rec['sku']
        x        = int(rec['nd']  / max_nd  * graph_w)
        y        = int(graph_h - rec['ros'] / max_ros * graph_h)
        img_path = sku_images_map.get(sku)
        placed   = _place_image(slide, img_path, x, y, img_h, cross_w=img_h,
                                rotate_tall=False)
        if not placed:
            not_found.add(sku)

    # Footnote: explain the frame and its relation to slide 1
    note_w = int(SLIDE_WIDTH - Cm(0.8))
    note_h = int(Cm(1.0))
    note_x = int(Cm(0.4))
    note_y = int(SLIDE_HEIGHT - note_h - Cm(0.2))
    txb = slide.shapes.add_textbox(note_x, note_y, note_w, note_h)
    tf = txb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = (
        "Рамка — альтернативное представление слайда 1: каждое изображение занимает "
        "ту же позицию, что и соответствующая точка на графике рассеяния "
        "(X = nd, Y = ros, начало координат — нижний левый угол). "
        "Весь диапазон данных равномерно масштабируется в размеры рамки."
    )
    run.font.name = FONT
    run.font.size = Pt(7)
    run.font.italic = True
    run.font.color.rgb = RGBColor(0x50, 0x50, 0x50)

    return not_found


def build_slide_three(prs, sku_data, thresholds, sku_images_map):
    """
    Slide 3: up to 4 side-by-side sections, one per category.
    Each section: [image column | data table].
    Columns at 1 cm, 6 pt font, only the rank cell is coloured.
    ros → 1 decimal; nd → percentage.
    Tall images (H > 2×W) are rotated 90° CCW (requires Pillow).
    """
    slide = prs.slides.add_slide(get_blank_layout(prs))

    sorted_data = sorted(sku_data, key=lambda r: r['rank'])
    if not sorted_data:
        return set()

    def cat_index(rank):
        for i, t in enumerate(thresholds):
            if rank <= t:
                return i
        return len(thresholds) if thresholds else 0

    n_cats = (len(thresholds) + 1) if thresholds else 1

    groups = [[] for _ in range(n_cats)]
    for rec in sorted_data:
        groups[min(cat_index(rec['rank']), n_cats - 1)].append(rec)

    active_groups = [(cat, grp) for cat, grp in enumerate(groups) if grp]
    if not active_groups:
        return set()

    # --- layout ---
    margin_t  = Cm(0.3)
    margin_b  = Cm(0.3)
    header_h  = Cm(0.9)
    col_w     = Cm(1.0)
    n_cols    = len(COLUMN_RENAME)
    table_w   = col_w * n_cols      # 5 cm
    img_col_w = Cm(1.5)
    gap_inner = Cm(0.1)             # image col ↔ table
    gap_outer = Cm(0.3)             # between sections

    n_active  = len(active_groups)
    section_w = img_col_w + gap_inner + table_w
    total_w   = n_active * section_w + (n_active - 1) * gap_outer
    start_x   = int((SLIDE_WIDTH - total_w) / 2)

    max_rows    = max(len(grp) for _, grp in active_groups)
    available_h = SLIDE_HEIGHT - margin_t - margin_b - header_h
    row_h       = max(Cm(0.5), min(Cm(1.2), int(available_h / max_rows)))

    col_order = list(COLUMN_RENAME.keys())
    not_found = set()

    for section_i, (cat, group) in enumerate(active_groups):
        color        = CATEGORY_COLORS[cat] if 0 <= cat < len(CATEGORY_COLORS) else None
        section_left = int(start_x + section_i * (section_w + gap_outer))
        table_left   = int(section_left + img_col_w + gap_inner)
        n_rows       = len(group)
        total_h      = int(header_h + row_h * n_rows)

        tbl_shape = slide.shapes.add_table(
            n_rows + 1, n_cols,
            table_left, int(margin_t),
            int(table_w), total_h,
        )
        tbl = tbl_shape.table

        for ci in range(n_cols):
            tbl.columns[ci].width = int(col_w)

        tbl.rows[0].height = int(header_h)
        for ci, key in enumerate(col_order):
            _set_cell_text(tbl.cell(0, ci), COLUMN_RENAME[key],
                           font_size=Pt(6), bold=True, center=True)

        for ri, rec in enumerate(group):
            row_idx = ri + 1
            tbl.rows[row_idx].height = int(row_h)

            for ci, key in enumerate(col_order):
                raw = rec.get(key, '')
                if key == 'ros' and raw is not None and raw != '':
                    display = f"{float(raw):.1f}"
                elif key == 'nd' and raw is not None and raw != '':
                    display = f"{float(raw) * 100:.0f}%"
                else:
                    display = raw

                cell = tbl.cell(row_idx, ci)
                _set_cell_text(cell, display, font_size=Pt(6), center=True)

                if ci == 0 and color:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = color

            sku      = rec['sku']
            img_top  = int(margin_t + header_h + ri * row_h)
            img_path = sku_images_map.get(sku)
            placed   = _place_image(
                slide, img_path,
                section_left, img_top, row_h, cross_w=img_col_w,
                rotate_tall=True,
            )
            if not placed:
                not_found.add(sku)

    return not_found


# ---------------------------------------------------------------------------
# Entry point
# ---------------------------------------------------------------------------
def main():
    try:
        print("=" * 50)
        print("  Запуск PPTX Velocity")
        print("=" * 50)

        if not os.path.exists(XLSX_PATH):
            write_error(
                "Файл velocity.xlsx не найден.\n"
                f"Пожалуйста, положите файл velocity.xlsx в папку example_folder "
                f"рядом с исполняемым файлом программы.\n"
                f"Ожидаемый путь: {XLSX_PATH}"
            )
            return

        # Ensure gray PNG exists and pre-process fallback image (white bg removal)
        _ensure_gray_png()
        _init_fallback()
        if os.path.exists(_FALLBACK_IMG):
            print(f"[OK] Изображение-заглушка: fallback.png")
        elif os.path.exists(_GENERATED_PNG):
            print(f"[OK] Изображение-заглушка: default_image.png (авто)")
        else:
            print(f"[!]  Изображение-заглушка недоступно — будет использоваться красный крестик")
        if not _HAS_PIL:
            print("[!]  Pillow не установлен — поворот вытянутых картинок на 3-м слайде отключён")

        print(f"\n[1/4] Чтение данных из {XLSX_PATH} ...")
        try:
            sku_data = read_sheet1()
        except (KeyError, ValueError) as e:
            write_error(str(e))
            return
        except Exception:
            write_error(
                "Не удалось прочитать данные из файла velocity.xlsx.\n"
                "Убедитесь, что:\n"
                "  1. Файл не открыт в Excel прямо сейчас.\n"
                "  2. Лист Sheet1 содержит корректную таблицу с заголовками.",
                technical_detail=traceback.format_exc(),
            )
            return

        if not sku_data:
            write_error(
                "Таблица на листе Sheet1 пустая — нет строк с данными. "
                "Добавьте данные под строкой заголовков."
            )
            return

        print(f"     Прочитано {len(sku_data)} SKU из Sheet1")

        graph_w, graph_h, thresholds = read_sheet2()
        sku_images_map = get_sku_images_map()
        print(f"     Найдено {len(sku_images_map)} изображений в папке example_folder")

        prs = Presentation()
        prs.slide_width  = SLIDE_WIDTH
        prs.slide_height = SLIDE_HEIGHT

        print("\n[2/4] Создание слайда 1 (точечный график) ...")
        build_slide_one(prs, sku_data)

        print("[3/4] Создание слайда 2 (карта изображений) ...")
        not_found_2 = build_slide_two(prs, sku_data, graph_w, graph_h, sku_images_map)

        print("[4/4] Создание слайда 3 (таблицы по категориям) ...")
        not_found_3 = build_slide_three(prs, sku_data, thresholds, sku_images_map)

        out_path = uniquify(OUTPUT_PATH)
        print(f"\nСохранение: {out_path} ...")
        prs.save(out_path)
        print(f"\n[ГОТОВО] Презентация сохранена:\n  {out_path}")

        all_missing = not_found_2 | not_found_3
        if all_missing:
            ext_list = ', '.join(IMAGE_EXTENSIONS)
            lines = [
                "Некоторые изображения не найдены в папке example_folder.\n"
                "Для каждого SKU ниже добавьте файл изображения в папку example_folder "
                f"рядом с программой (поддерживаемые форматы: {ext_list}):\n\n"
            ]
            for sku in sorted(all_missing):
                lines.append(f"  - '{sku}': файл должен называться '{sku}.jpg' "
                              f"(или другой поддерживаемый формат).\n")
            msg = "".join(lines)
            txt_path = uniquify(out_path.replace('.pptx', '.txt'))
            with open(txt_path, 'w', encoding='utf-8') as f:
                f.write(msg)
            print(f"\n[!] {msg}")

        print("\nНажмите Enter для закрытия...")
        input()

    except Exception:
        write_error(
            "Произошла непредвиденная ошибка при создании презентации.\n"
            "Убедитесь, что файл velocity.xlsx корректен и не открыт в Excel.",
            technical_detail=traceback.format_exc(),
        )
        print("\nНажмите Enter для закрытия...")
        input()


if __name__ == "__main__":
    main()
